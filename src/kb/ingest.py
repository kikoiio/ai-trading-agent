"""PPTX ingestion pipeline: parse slides, extract images, describe charts via vision model."""

import base64
import io
import json
import logging
import os
import re
import time
from pathlib import Path

import requests
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from rich.console import Console
from rich.progress import Progress, SpinnerColumn, TextColumn, BarColumn, MofNCompleteColumn

from src.config_loader import CONFIG
from src.kb.models import KBEntry, auto_tag
from src.kb.vectorstore import KBVectorStore

logger = logging.getLogger(__name__)
console = Console()


def _extract_chapter_number(folder_name: str) -> str:
    """Extract chapter number from folder name like '05_Trends' or 'Chapter_05'."""
    match = re.search(r"(\d+)", folder_name)
    return match.group(1).zfill(2) if match else "00"


def _extract_slide_text(slide) -> tuple[str, str]:
    """Extract title and body text from a slide.

    Returns (title, body_text).
    """
    title = ""
    body_parts = []

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue
        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and hasattr(shape, "placeholder_format"):
            # idx 0 is typically the title placeholder
            if shape.placeholder_format.idx == 0:
                title = text
                continue
        body_parts.append(text)

    if not title and body_parts:
        title = body_parts.pop(0)

    return title, "\n".join(body_parts)


def _extract_slide_images(slide, output_dir: str, entry_id: str) -> list[str]:
    """Extract images from a slide, save to output_dir, return list of paths."""
    os.makedirs(output_dir, exist_ok=True)
    image_paths = []
    img_index = 0

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            img_index += 1
            image = shape.image
            ext = image.content_type.split("/")[-1] if image.content_type else "png"
            if ext == "jpeg":
                ext = "jpg"
            filename = f"{entry_id}_img{img_index}.{ext}"
            filepath = os.path.join(output_dir, filename)

            with open(filepath, "wb") as f:
                f.write(image.blob)

            image_paths.append(filepath)

    return image_paths


def _describe_image_with_vision(image_path: str, vision_model: str) -> str:
    """Send an image to a vision model via OpenRouter and get a chart description."""
    api_key = CONFIG["openrouter_api_key"]
    base_url = CONFIG["openrouter_base_url"]

    # Read and base64 encode the image
    try:
        with Image.open(image_path) as img:
            # Convert to RGB PNG if needed (handles EMF, WMF, etc.)
            if img.mode not in ("RGB", "RGBA"):
                img = img.convert("RGB")
            buf = io.BytesIO()
            img.save(buf, format="PNG")
            b64 = base64.b64encode(buf.getvalue()).decode("utf-8")
    except Exception as e:
        logger.warning("Failed to read image %s: %s", image_path, e)
        return f"[Image could not be processed: {e}]"

    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json",
    }

    payload = {
        "model": vision_model,
        "messages": [
            {
                "role": "user",
                "content": [
                    {
                        "type": "text",
                        "text": (
                            "Describe this trading chart image in detail for a knowledge base. Focus on:\n"
                            "- Price action patterns visible (bars, candles, trends, channels, trading ranges)\n"
                            "- Any annotations, arrows, labels, or text drawn on the chart\n"
                            "- The specific trading concept being illustrated (Al Brooks style price action)\n"
                            "- Entry/exit points or signal bars if marked\n"
                            "- The overall market context (trend direction, strength, reversals)\n"
                            "Keep the description factual, detailed, and searchable. "
                            "Use Al Brooks terminology where applicable."
                        ),
                    },
                    {
                        "type": "image_url",
                        "image_url": {
                            "url": f"data:image/png;base64,{b64}",
                        },
                    },
                ],
            }
        ],
        "max_tokens": 500,
        "temperature": 0.2,
    }

    try:
        resp = requests.post(
            f"{base_url}/chat/completions",
            headers=headers,
            json=payload,
            timeout=60,
        )
        resp.raise_for_status()
        content = resp.json()["choices"][0]["message"]["content"]
        return content.strip()
    except Exception as e:
        logger.error("Vision API call failed for %s: %s", image_path, e)
        return f"[Vision description failed: {e}]"


def ingest_pptx_folder(
    source_dir: str,
    data_dir: str = "data/kb",
    vision_model: str | None = None,
    skip_images: bool = False,
    reset: bool = False,
    delay: float = 1.0,
    chapters_filter: list[str] | None = None,
) -> dict:
    """Ingest all PPTX files from the source directory into the knowledge base.

    Args:
        source_dir: Path to Al_Brooks_trading_course/ folder.
        data_dir: Path to KB data directory.
        vision_model: Vision model for chart descriptions (default from config).
        skip_images: If True, skip vision processing of images.
        reset: If True, clear existing KB before ingesting.
        delay: Seconds between vision API calls.
        chapters_filter: Optional list of chapter numbers to process (e.g., ["01", "05"]).

    Returns:
        Dict with ingestion statistics.
    """
    source_path = Path(source_dir)
    if not source_path.exists():
        raise FileNotFoundError(f"Source directory not found: {source_dir}")

    vision_model = vision_model or CONFIG.get("kb_vision_model") or "openai/gpt-4o-mini"
    images_dir = os.path.join(data_dir, "images")
    os.makedirs(images_dir, exist_ok=True)

    store = KBVectorStore(persist_dir=os.path.join(data_dir, "chroma_db"))
    if reset:
        store.delete_all()
        console.print("[yellow]Cleared existing knowledge base[/yellow]")

    # Discover chapters (subfolders)
    chapter_dirs = sorted([
        d for d in source_path.iterdir()
        if d.is_dir() and not d.name.startswith(".")
    ])

    if chapters_filter:
        chapter_dirs = [
            d for d in chapter_dirs
            if _extract_chapter_number(d.name) in chapters_filter
        ]

    stats = {
        "chapters_processed": 0,
        "slides_processed": 0,
        "images_described": 0,
        "entries_stored": 0,
        "errors": [],
    }

    with Progress(
        SpinnerColumn(),
        TextColumn("[progress.description]{task.description}"),
        BarColumn(),
        MofNCompleteColumn(),
        console=console,
    ) as progress:
        chapter_task = progress.add_task("Chapters", total=len(chapter_dirs))

        for chapter_dir in chapter_dirs:
            chapter_name = chapter_dir.name
            chapter_num = _extract_chapter_number(chapter_name)
            chapter_images_dir = os.path.join(images_dir, chapter_name)

            if chapters_filter and chapter_num not in chapters_filter:
                progress.advance(chapter_task)
                continue

            # Find all PPTX files in this chapter
            pptx_files = sorted(chapter_dir.glob("*.pptx"))
            if not pptx_files:
                # Also check for files directly (not in subfolder)
                pptx_files = sorted(chapter_dir.glob("**/*.pptx"))

            progress.update(chapter_task, description=f"Chapter {chapter_num}: {chapter_name}")

            for pptx_file in pptx_files:
                try:
                    prs = Presentation(str(pptx_file))
                except Exception as e:
                    error_msg = f"Failed to open {pptx_file}: {e}"
                    logger.error(error_msg)
                    stats["errors"].append(error_msg)
                    continue

                relative_path = str(pptx_file.relative_to(source_path))
                slide_task = progress.add_task(
                    f"  {pptx_file.name}",
                    total=len(prs.slides),
                )

                batch_entries = []

                for slide_idx, slide in enumerate(prs.slides, start=1):
                    entry_id = f"AB-CH{chapter_num}-S{slide_idx:02d}"

                    # Extract text
                    title, body = _extract_slide_text(slide)
                    if not title:
                        title = f"Slide {slide_idx}"

                    # Extract and describe images
                    image_paths = []
                    image_descriptions = []

                    if not skip_images:
                        image_paths = _extract_slide_images(
                            slide, chapter_images_dir, entry_id
                        )
                        for img_path in image_paths:
                            progress.update(slide_task, description=f"  Describing {os.path.basename(img_path)}")
                            description = _describe_image_with_vision(img_path, vision_model)
                            image_descriptions.append(description)
                            stats["images_described"] += 1
                            if delay > 0:
                                time.sleep(delay)

                    # Build KB entry
                    entry = KBEntry(
                        id=entry_id,
                        source_file=relative_path,
                        chapter=chapter_name,
                        slide_number=slide_idx,
                        slide_title=title,
                        text_content=body,
                        image_descriptions=image_descriptions,
                        tags=auto_tag(f"{title} {body} {' '.join(image_descriptions)}"),
                        image_paths=[os.path.relpath(p, data_dir) for p in image_paths],
                    )
                    entry.build_combined_text()
                    batch_entries.append(entry)
                    stats["slides_processed"] += 1
                    progress.advance(slide_task)

                # Store batch for this PPTX file
                if batch_entries:
                    store.add_entries_batch(batch_entries)
                    stats["entries_stored"] += len(batch_entries)

                progress.remove_task(slide_task)

            stats["chapters_processed"] += 1
            progress.advance(chapter_task)

    # Print summary
    console.print(f"\n[green]Ingestion complete![/green]")
    console.print(f"  Chapters: {stats['chapters_processed']}")
    console.print(f"  Slides: {stats['slides_processed']}")
    console.print(f"  Images described: {stats['images_described']}")
    console.print(f"  KB entries stored: {stats['entries_stored']}")
    if stats["errors"]:
        console.print(f"  [red]Errors: {len(stats['errors'])}[/red]")
        for err in stats["errors"]:
            console.print(f"    - {err}")

    return stats
